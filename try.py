import os
from tkinter import Tk, filedialog
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    """Extract text from a single PPTX file."""
    prs = Presentation(pptx_path)
    text_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_content.append(f"\n--- Slide {slide_num} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return "\n".join(text_content)

def main():
    # Hide the root Tkinter window
    root = Tk()
    root.withdraw()

    # Let user select multiple PPTX files
    pptx_files = filedialog.askopenfilenames(
        title="Select PowerPoint Files",
        filetypes=[("PowerPoint files", "*.pptx")]
    )

    if not pptx_files:
        print("No files selected.")
        return

    # Create an output folder in the same directory as the script
    output_folder = os.path.join(os.getcwd(), "pptx_to_txt_output")
    os.makedirs(output_folder, exist_ok=True)

    # Process each selected PPTX file
    for pptx_path in pptx_files:
        text = extract_text_from_pptx(pptx_path)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        output_path = os.path.join(output_folder, f"{base_name}.txt")

        with open(output_path, "w", encoding="utf-8") as txt_file:
            txt_file.write(text)

        print(f"✅ Converted: {base_name}.pptx → {base_name}.txt")

    print(f"\nAll text files are saved in: {output_folder}")

if __name__ == "__main__":
    main()
